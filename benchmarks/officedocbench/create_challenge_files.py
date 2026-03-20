#!/usr/bin/env python3
"""Create challenge test files that expose DocParse's known parsing gaps.

These are small files designed to stress features that aren't fully supported yet.
They serve as stretch goals for future parser improvements.

Usage:
    uv run benchmarks/officedocbench/create_challenge_files.py
"""

from __future__ import annotations

import json
from pathlib import Path

CHALLENGE_DIR = Path(__file__).parent.parent.parent / "data" / "test_files" / "challenge"
GT_DIR = Path(__file__).parent / "ground_truth"


def create_docx_footnotes():
    """DOCX with real footnote and endnote content (not just separators)."""
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        print("  SKIP challenge_footnotes.docx (python-docx not installed)")
        return None

    doc = Document()
    doc.add_heading("Document with Footnotes", level=1)
    p = doc.add_paragraph("This paragraph has a footnote")

    # Add footnote reference
    run = p.runs[0] if p.runs else p.add_run("")
    footnote_ref = OxmlElement("w:r")
    rpr = OxmlElement("w:rPr")
    style = OxmlElement("w:rStyle")
    style.set(qn("w:val"), "FootnoteReference")
    rpr.append(style)
    footnote_ref.append(rpr)
    sup = OxmlElement("w:footnoteReference")
    sup.set(qn("w:id"), "1")
    footnote_ref.append(sup)
    run._element.addnext(footnote_ref)

    doc.add_paragraph("Second paragraph with an endnote reference.")
    doc.add_heading("Conclusion", level=2)
    doc.add_paragraph("Final thoughts on the topic.")

    path = CHALLENGE_DIR / "challenge_footnotes.docx"
    doc.save(str(path))

    return {
        "file": "challenge/challenge_footnotes.docx",
        "format": "docx",
        "source": "generated (python-docx)",
        "verified": True,
        "features": {
            "text": {"paragraph_count": 3, "total_words": 20, "key_phrases": ["Document with Footnotes"]},
            "headings": {"present": True, "count": 2, "by_level": {"1": 1, "2": 1}},
            "tables": {"present": False, "count": 0, "total_rows": 0, "has_merged_cells": False},
            "track_changes": {"present": False, "count": 0, "types": {}, "authors": []},
            "comments": {"present": False, "count": 0, "authors": [], "texts": []},
            "headers_footers": {"present": False, "header_count": 0, "footer_count": 0},
            "footnotes_endnotes": {"present": True, "count": 1},
            "text_boxes": {"present": False, "count": 0},
            "images": {"present": False, "count": 0},
            "lists": {"present": False, "count": 0},
            "metadata": {"title": "", "author": "", "created": "", "modified": ""},
            "sheets": None,
        },
        "full_text_words": ["conclusion", "document", "endnote", "final", "footnote", "footnotes",
                            "has", "on", "paragraph", "reference", "second", "the", "this",
                            "thoughts", "topic", "with"],
    }


def create_pptx_speaker_notes():
    """PPTX with speaker notes on slides."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("  SKIP challenge_speaker_notes.pptx (python-pptx not installed)")
        return None

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content

    # Slide 1 with speaker notes
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Introduction"
    slide1.placeholders[1].text = "Welcome to the presentation"
    notes_slide = slide1.notes_slide
    notes_slide.notes_text_frame.text = "Remember to greet the audience warmly"

    # Slide 2 with speaker notes
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Key Points"
    slide2.placeholders[1].text = "Three important items"
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "Spend 5 minutes on this slide. Emphasize point two."

    # Slide 3 without notes
    slide3 = prs.slides.add_slide(slide_layout)
    slide3.shapes.title.text = "Thank You"
    slide3.placeholders[1].text = "Questions?"

    path = CHALLENGE_DIR / "challenge_speaker_notes.pptx"
    prs.save(str(path))

    return {
        "file": "challenge/challenge_speaker_notes.pptx",
        "format": "pptx",
        "source": "generated (python-pptx)",
        "verified": True,
        "features": {
            "text": {"paragraph_count": 3, "total_words": 25, "key_phrases": ["Introduction", "Key Points"]},
            "headings": {"present": True, "count": 3, "by_level": {"1": 3}},
            "tables": None,
            "track_changes": None,
            "comments": None,
            "headers_footers": None,
            "footnotes_endnotes": None,
            "speaker_notes": {"present": True, "count": 2},
            "text_boxes": {"present": False, "count": 0},
            "images": {"present": False, "count": 0},
            "lists": {"present": False, "count": 0},
            "metadata": {"title": "", "author": "", "created": "", "modified": ""},
            "sheets": None,
        },
        "full_text_words": ["5", "audience", "emphasize", "greet", "important", "introduction",
                            "items", "key", "minutes", "on", "point", "points", "presentation",
                            "questions", "remember", "slide", "spend", "thank", "the", "this",
                            "three", "to", "two", "warmly", "welcome", "you"],
    }


def create_xlsx_formulas():
    """XLSX with formulas and multiple sheets."""
    try:
        from openpyxl import Workbook
    except ImportError:
        print("  SKIP challenge_formulas.xlsx (openpyxl not installed)")
        return None

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Revenue"
    ws1["A1"] = "Quarter"
    ws1["B1"] = "Revenue"
    ws1["C1"] = "Expenses"
    ws1["D1"] = "Profit"
    ws1["A2"] = "Q1"
    ws1["B2"] = 50000
    ws1["C2"] = 30000
    ws1["D2"] = "=B2-C2"
    ws1["A3"] = "Q2"
    ws1["B3"] = 65000
    ws1["C3"] = 35000
    ws1["D3"] = "=B3-C3"
    ws1["A4"] = "Total"
    ws1["B4"] = "=SUM(B2:B3)"
    ws1["C4"] = "=SUM(C2:C3)"
    ws1["D4"] = "=SUM(D2:D3)"

    # Merge cells for a header
    ws1.merge_cells("A6:D6")
    ws1["A6"] = "Financial Summary FY2026"

    ws2 = wb.create_sheet("Employees")
    ws2["A1"] = "Name"
    ws2["B1"] = "Department"
    ws2["A2"] = "Alice"
    ws2["B2"] = "Engineering"
    ws2["A3"] = "Bob"
    ws2["B3"] = "Sales"

    path = CHALLENGE_DIR / "challenge_formulas.xlsx"
    wb.save(str(path))

    return {
        "file": "challenge/challenge_formulas.xlsx",
        "format": "xlsx",
        "source": "generated (openpyxl)",
        "verified": True,
        "features": {
            "text": {"paragraph_count": 0, "total_words": 20, "key_phrases": ["Financial Summary FY2026"]},
            "headings": None,
            "tables": {"present": True, "count": 2, "total_rows": 6, "has_merged_cells": True},
            "track_changes": None,
            "comments": None,
            "headers_footers": None,
            "footnotes_endnotes": None,
            "speaker_notes": None,
            "text_boxes": None,
            "images": {"present": False, "count": 0},
            "metadata": {"title": "", "author": "", "created": "", "modified": ""},
            "sheets": {"present": True, "names": ["Revenue", "Employees"]},
            "lists": None,
        },
        "full_text_words": ["alice", "bob", "department", "employees", "engineering",
                            "expenses", "financial", "fy2026", "name", "profit",
                            "q1", "q2", "quarter", "revenue", "sales", "summary", "total"],
    }


def create_docx_nested_lists():
    """DOCX with deeply nested multi-level lists."""
    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError:
        print("  SKIP challenge_nested_lists.docx (python-docx not installed)")
        return None

    doc = Document()
    doc.add_heading("Project Requirements", level=1)

    # Flat list items (python-docx doesn't natively support nested lists easily,
    # so we use indented paragraphs)
    doc.add_paragraph("Frontend", style="List Bullet")
    doc.add_paragraph("React components", style="List Bullet 2")
    doc.add_paragraph("Button component", style="List Bullet 3")
    doc.add_paragraph("Form component", style="List Bullet 3")
    doc.add_paragraph("State management", style="List Bullet 2")
    doc.add_paragraph("Backend", style="List Bullet")
    doc.add_paragraph("API endpoints", style="List Bullet 2")
    doc.add_paragraph("Database", style="List Bullet 2")
    doc.add_paragraph("PostgreSQL schema", style="List Bullet 3")
    doc.add_paragraph("Migrations", style="List Bullet 3")

    doc.add_heading("Timeline", level=2)
    doc.add_paragraph("Phase 1: Design (2 weeks)")
    doc.add_paragraph("Phase 2: Implementation (4 weeks)")

    path = CHALLENGE_DIR / "challenge_nested_lists.docx"
    doc.save(str(path))

    return {
        "file": "challenge/challenge_nested_lists.docx",
        "format": "docx",
        "source": "generated (python-docx)",
        "verified": True,
        "features": {
            "text": {"paragraph_count": 2, "total_words": 25,
                     "key_phrases": ["Project Requirements", "Frontend"]},
            "headings": {"present": True, "count": 2, "by_level": {"1": 1, "2": 1}},
            "tables": {"present": False, "count": 0, "total_rows": 0, "has_merged_cells": False},
            "track_changes": {"present": False, "count": 0, "types": {}, "authors": []},
            "comments": {"present": False, "count": 0, "authors": [], "texts": []},
            "headers_footers": {"present": False, "header_count": 0, "footer_count": 0},
            "footnotes_endnotes": {"present": False, "count": 0},
            "text_boxes": {"present": False, "count": 0},
            "images": {"present": False, "count": 0},
            "lists": {"present": True, "count": 2},
            "metadata": {"title": "", "author": "", "created": "", "modified": ""},
            "sheets": None,
        },
        "full_text_words": ["api", "backend", "button", "component", "components", "database",
                            "design", "endpoints", "form", "frontend", "implementation",
                            "management", "migrations", "phase", "postgresql", "project",
                            "react", "requirements", "schema", "state", "timeline", "weeks"],
    }


def create_html_complex_structure():
    """HTML with tables, nested lists, images, and semantic elements."""
    html_content = """<!DOCTYPE html>
<html>
<head>
<title>Complex HTML Document</title>
<meta name="author" content="Test Author">
</head>
<body>
<h1>Report Title</h1>
<p>Executive summary paragraph.</p>

<h2>Data Section</h2>
<table>
<thead>
<tr><th>Metric</th><th>Q1</th><th>Q2</th></tr>
</thead>
<tbody>
<tr><td>Revenue</td><td>$50K</td><td>$65K</td></tr>
<tr><td>Users</td><td>1,000</td><td>1,500</td></tr>
</tbody>
</table>

<h2>Action Items</h2>
<ol>
<li>Review Q1 numbers
  <ul>
    <li>Check revenue figures</li>
    <li>Verify user counts</li>
  </ul>
</li>
<li>Plan Q3 strategy</li>
</ol>

<h3>Notes</h3>
<p>Important: all figures are preliminary.</p>

<figure>
<img src="chart.png" alt="Revenue chart showing growth trend">
<figcaption>Figure 1: Revenue growth</figcaption>
</figure>
</body>
</html>"""

    path = CHALLENGE_DIR / "challenge_complex.html"
    path.write_text(html_content)

    return {
        "file": "challenge/challenge_complex.html",
        "format": "html",
        "source": "generated",
        "verified": True,
        "features": {
            "text": {"paragraph_count": 3, "total_words": 35,
                     "key_phrases": ["Report Title", "Executive summary paragraph"]},
            "headings": {"present": True, "count": 3, "by_level": {"1": 1, "2": 2, "3": 1}},
            "tables": {"present": True, "count": 1, "total_rows": 2, "has_merged_cells": False},
            "track_changes": None,
            "comments": None,
            "headers_footers": None,
            "footnotes_endnotes": None,
            "speaker_notes": None,
            "text_boxes": None,
            "images": {"present": True, "count": 1},
            "lists": {"present": True, "count": 2},
            "metadata": {"title": "Complex HTML Document", "author": "Test Author", "created": "", "modified": ""},
            "sheets": None,
        },
        "full_text_words": ["1", "50k", "65k", "action", "all", "are", "chart", "check",
                            "complex", "counts", "data", "document", "executive", "figures",
                            "growth", "html", "important", "items", "metric", "notes", "numbers",
                            "paragraph", "plan", "preliminary", "q1", "q2", "q3", "report",
                            "revenue", "review", "section", "showing", "strategy", "summary",
                            "title", "trend", "user", "users", "verify"],
    }


def main():
    CHALLENGE_DIR.mkdir(parents=True, exist_ok=True)

    generators = [
        ("challenge_footnotes.docx", create_docx_footnotes),
        ("challenge_speaker_notes.pptx", create_pptx_speaker_notes),
        ("challenge_formulas.xlsx", create_xlsx_formulas),
        ("challenge_nested_lists.docx", create_docx_nested_lists),
        ("challenge_complex.html", create_html_complex_structure),
    ]

    created = 0
    for name, gen_fn in generators:
        print(f"  Creating {name}...", end=" ")
        gt = gen_fn()
        if gt is None:
            continue

        # Save ground truth
        gt_path = GT_DIR / f"{gt['file'].replace('/', '_').replace('challenge_', 'challenge_')}.json"
        # Use simpler naming: challenge_X.ext.json
        gt_filename = name + ".json"
        gt_path = GT_DIR / gt_filename
        with open(gt_path, "w") as f:
            json.dump(gt, f, indent=2)

        created += 1
        print("OK")

    print(f"\nCreated {created} challenge files in {CHALLENGE_DIR}")
    print(f"Ground truth in {GT_DIR}")
    print("\nThese files expose known gaps:")
    print("  - challenge_footnotes.docx: real footnote content (not just separators)")
    print("  - challenge_speaker_notes.pptx: speaker notes on slides")
    print("  - challenge_formulas.xlsx: formulas + merged cells + multi-sheet")
    print("  - challenge_nested_lists.docx: deeply nested bullet lists")
    print("  - challenge_complex.html: nested lists + figure/img + semantic elements")


if __name__ == "__main__":
    main()
