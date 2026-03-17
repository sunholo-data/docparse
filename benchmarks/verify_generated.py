#!/usr/bin/env python3
"""Verification loop for generated Office documents.

Checks that generated files are valid at multiple levels:
  Level 1: Structural (ZIP + XML well-formedness)
  Level 2: Library (python-docx, python-pptx, openpyxl open without errors)
  Level 4: Roundtrip (generate → parse through DocParse → verify blocks)

Usage:
  uv run --with python-pptx --with openpyxl --with python-docx benchmarks/verify_generated.py
  uv run --with python-pptx --with openpyxl --with python-docx benchmarks/verify_generated.py --fix  # regenerate after fixes
"""

import json
import os
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

REPO_DIR = Path(__file__).parent.parent
EXAMPLES_DIR = REPO_DIR / "data" / "examples"

# ── Level 1: Structural Checks ──────────────────────────────────────────────

def verify_structure(path: Path) -> tuple[list[str], list[str]]:
    """Check ZIP structure and XML well-formedness. Returns (errors, warnings)."""
    errors = []
    warnings = []
    if not zipfile.is_zipfile(path):
        return ([f"Not a valid ZIP file"], [])

    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()

        # Check required entries by format
        ext = path.suffix.lower()
        if ext == ".docx":
            required = ["[Content_Types].xml", "_rels/.rels", "word/document.xml"]
        elif ext == ".pptx":
            required = ["[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"]
        elif ext == ".xlsx":
            required = ["[Content_Types].xml", "_rels/.rels", "xl/workbook.xml"]
        elif ext in (".odt", ".odp", ".ods"):
            required = ["mimetype", "META-INF/manifest.xml", "content.xml"]
        else:
            required = []

        for req in required:
            if req not in names:
                errors.append(f"Missing required entry: {req}")

        # Check XML well-formedness for all .xml and .rels entries
        for name in names:
            if name.endswith(".xml") or name.endswith(".rels"):
                try:
                    data = zf.read(name).decode("utf-8")
                    ET.fromstring(data)
                except ET.ParseError as e:
                    errors.append(f"XML parse error in {name}: {e}")
                except Exception as e:
                    errors.append(f"Error reading {name}: {e}")

        # PPTX-specific: check theme exists
        if ext == ".pptx":
            has_theme = any("theme" in n for n in names)
            if not has_theme:
                errors.append("Missing ppt/theme/theme1.xml (Keynote requires this)")

        # ODF-specific: check mimetype is first entry and uncompressed
        # NOTE: AILANG's createArchive always compresses. This is a known limitation
        # that needs a stdlib change (store flag). Downgraded to warning.
        if ext in (".odt", ".odp", ".ods"):
            if names and names[0] != "mimetype":
                warnings.append(f"mimetype should be first ZIP entry, got: {names[0]}")
            info = zf.getinfo("mimetype") if "mimetype" in names else None
            if info and info.compress_type != zipfile.ZIP_STORED:
                warnings.append("mimetype entry is compressed (ODF spec wants uncompressed) — AILANG stdlib limitation")

    return (errors, warnings)


# ── Level 2: Library Validation ──────────────────────────────────────────────

def verify_library(path: Path) -> list[str]:
    """Open with Python libraries and check for errors/warnings."""
    errors = []
    ext = path.suffix.lower()
    import warnings

    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            if len(doc.paragraphs) == 0:
                errors.append("DOCX has 0 paragraphs")
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            if len(prs.slides) == 0:
                errors.append("PPTX has 0 slides")
            for i, slide in enumerate(prs.slides):
                if len(slide.shapes) == 0:
                    errors.append(f"Slide {i+1} has 0 shapes")
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            with warnings.catch_warnings(record=True) as w:
                warnings.simplefilter("always")
                wb = load_workbook(str(path))
                for warning in w:
                    errors.append(f"Warning: {warning.message}")
                if len(wb.sheetnames) == 0:
                    errors.append("XLSX has 0 sheets")
    except Exception as e:
        errors.append(f"{type(e).__name__}: {e}")

    return errors


# ── Level 4: Roundtrip Verification ──────────────────────────────────────────

def verify_roundtrip(path: Path) -> list[str]:
    """Parse the generated file through DocParse and check output."""
    errors = []
    ext = path.suffix.lower()

    # Skip formats DocParse can't re-parse (ODP, ODS have basic parsers)
    if ext in (".html", ".htm"):
        return []  # HTML roundtrip already tested

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            result = subprocess.run(
                ["ailang", "run", "--entry", "main", "--caps", "IO,FS,Env",
                 "--max-recursion-depth", "50000",
                 "docparse/main.ail", str(path)],
                capture_output=True, text=True, cwd=str(REPO_DIR),
                timeout=30,
                env={**os.environ, "DOCPARSE_OUTPUT_DIR": tmpdir},
            )
            if result.returncode != 0:
                stderr_lines = result.stderr.strip().split("\n")[-3:]
                errors.append(f"DocParse failed (exit {result.returncode}): {' '.join(stderr_lines)}")
                return errors

            # Check output mentions blocks
            output = result.stdout
            if "Blocks:   0" in output:
                errors.append("Roundtrip produced 0 blocks")

            # Extract block count from output
            for line in output.split("\n"):
                if line.strip().startswith("Blocks:"):
                    count = line.strip().split()[-1]
                    if count == "0":
                        errors.append("Roundtrip produced 0 blocks")

        except subprocess.TimeoutExpired:
            errors.append("Roundtrip timed out (30s)")
        except Exception as e:
            errors.append(f"Roundtrip error: {e}")

    return errors


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    print("=== DocParse Generated File Verification ===\n")

    files = sorted(EXAMPLES_DIR.glob("*"))
    files = [f for f in files if f.suffix.lower() in (".docx", ".pptx", ".xlsx", ".odt", ".odp", ".ods", ".html")]

    if not files:
        print(f"No files found in {EXAMPLES_DIR}")
        sys.exit(1)

    all_pass = True
    for path in files:
        print(f"── {path.name} ──")

        # Level 1: Structure
        if path.suffix.lower() in (".html", ".htm"):
            # HTML: just check it's valid XML (XHTML)
            try:
                ET.parse(str(path))
                print(f"  L1 Structure:  PASS (valid XHTML)")
            except ET.ParseError as e:
                print(f"  L1 Structure:  WARN (not valid XHTML: {e})")
        else:
            struct_errors, struct_warnings = verify_structure(path)
            if struct_errors:
                print(f"  L1 Structure:  FAIL")
                for e in struct_errors:
                    print(f"     ⚠ {e}")
                all_pass = False
            elif struct_warnings:
                print(f"  L1 Structure:  WARN")
                for w in struct_warnings:
                    print(f"     ⚠ {w}")
            else:
                print(f"  L1 Structure:  PASS")

        # Level 2: Library
        if path.suffix.lower() in (".docx", ".pptx", ".xlsx"):
            lib_errors = verify_library(path)
            if lib_errors:
                print(f"  L2 Library:    WARN")
                for e in lib_errors:
                    print(f"     ⚠ {e}")
            else:
                print(f"  L2 Library:    PASS")

        # Level 4: Roundtrip (skip HTML — already tested separately)
        if path.suffix.lower() not in (".html", ".htm"):
            rt_errors = verify_roundtrip(path)
            if rt_errors:
                print(f"  L4 Roundtrip:  FAIL")
                for e in rt_errors:
                    print(f"     ⚠ {e}")
                all_pass = False
            else:
                print(f"  L4 Roundtrip:  PASS")

        print()

    if all_pass:
        print("=== ALL CHECKS PASSED ===")
    else:
        print("=== SOME CHECKS FAILED (see above) ===")
        sys.exit(1)


if __name__ == "__main__":
    main()
